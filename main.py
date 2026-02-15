import os
import re
from io import BytesIO
import asyncio

from telegram import (
    Update,
    InlineKeyboardMarkup,
    InlineKeyboardButton,
    InputFile,
)
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    ContextTypes,
    filters,
)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor

# ================== НАЛАШТУВАННЯ ==================

import os
from dotenv import load_dotenv

load_dotenv()  # локально прочитає файл .env

BOT_TOKEN = os.getenv("BOT_TOKEN")
if not BOT_TOKEN:
    raise RuntimeError("BOT_TOKEN не знайдено! Додай його в .env або в змінні середовища.")


# Розміри слайда (16:9)
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5

# Область для тексту (1 см ~ 0.39 inch відступ з усіх боків)
TEXT_AREA_LEFT_IN = 0.4
TEXT_AREA_TOP_IN = 0.4
TEXT_AREA_WIDTH_IN = SLIDE_WIDTH_IN - 2 * TEXT_AREA_LEFT_IN
TEXT_AREA_HEIGHT_IN = SLIDE_HEIGHT_IN - 2 * TEXT_AREA_TOP_IN

# Шрифти
AVAILABLE_FONTS = ["Arial Black", "Arial", "Times New Roman", "Calibri"]
AVAILABLE_FONT_SIZES = [38, 40, 42, 44, 46, 48]

DEFAULT_SETTINGS = {
    "font_name": "Arial Black",
    "font_size": 48,
    "bg_mode": "none",  # none | color | image
    "bg_color": (255, 255, 255),
    "bg_image": None,   # BytesIO
}

# Час очікування між повідомленнями (в секундах)
TEXT_ACCUMULATION_TIMEOUT = 3.0

def get_max_chars_for_font_size(font_size: int) -> int:
    """Орієнтовний ліміт символів для одного слайду залежно від розміру шрифту."""
    if font_size >= 48:
        return 130
    elif font_size >= 46:
        return 140
    elif font_size >= 44:
        return 155
    elif font_size >= 42:
        return 170
    elif font_size >= 40:
        return 185
    else:  # 38 і менше
        return 200


# ================== ЛОГІКА ДЛЯ ЗВИЧАЙНОГО ТЕКСТУ ==================

def split_by_parentheses(line: str):
    """Розбиває рядок на звичайний текст і блоки в () як окремі шматки."""
    parts = []
    pos = 0
    for m in re.finditer(r"\([^)]*\)", line):
        start, end = m.span()
        if start > pos:
            before = line[pos:start].strip()
            if before:
                parts.append(before)
        par = line[start:end].strip()
        if par:
            parts.append(par)
        pos = end
    if pos < len(line):
        tail = line[pos:].strip()
        if tail:
            parts.append(tail)
    return parts

def sentence_split(text: str):
    """Ділить текст на речення за . ? !, не рве слова."""
    parts = re.split(r"([.!?])", text)
    sentences = []
    cur = ""
    for p in parts:
        if not p:
            continue
        cur += p
        if p in ".!?":
            sentences.append(cur.strip())
            cur = ""
    if cur.strip():
        sentences.append(cur.strip())
    return sentences

def chunk_text_for_slides(text: str, max_chars: int):
    """Ділить текст на шматки по реченнях до max_chars."""
    sentences = sentence_split(text)
    if not sentences:
        return [text]

    chunks = []
    cur = ""

    for s in sentences:
        if not cur:
            cur = s
            continue

        candidate = cur + " " + s
        if len(candidate) <= max_chars:
            cur = candidate
        else:
            chunks.append(cur)
            cur = s

    if cur:
        chunks.append(cur)

    return chunks


def text_to_chunks(raw_text: str, settings) -> list[str]:
    """
    Простий режим:
    - кожен непорожній рядок = репліка
    - ( ... ) → окремі шматки
    - далі по реченнях до ліміту, який залежить від розміру шрифту
    """
    max_chars = get_max_chars_for_font_size(settings["font_size"])
    lines = [l.strip() for l in raw_text.splitlines()]
    chunks: list[str] = []

    for line in lines:
        if not line:
            continue
        parts = split_by_parentheses(line)
        for part in parts:
            for c in chunk_text_for_slides(part, max_chars=max_chars):
                chunks.append(c)

    return chunks


# ================== ЛОГІКА СЦЕНАРІЮ ЗІ СПИСОК ДІЙОВИХ ОСІБ ==================

def parse_replicas_from_colon_headers(raw_text: str):
    """
    Формат:
    Sternsinger 2:
    текст...

    (ремарка між репліками)

    Engel 1:
    ...
    """
    lines = [l.rstrip("\n") for l in raw_text.splitlines()]

    current_speaker: str | None = None
    current_buffer: list[str] = []   # рядки всередині репліки
    replicas: list[tuple[str, str]] = []
    remark_blocks: list[str] = []    # окремі ремарки між репліками

    def flush_replica():
        nonlocal current_speaker, current_buffer, replicas
        text = "\n".join(current_buffer).strip()
        if current_speaker and text:
            replicas.append((current_speaker, text))
        current_buffer = []

    for line in lines:
        stripped = line.strip()

        # новий персонаж: рядок, що закінчується на ":"
        if stripped.endswith(":") and len(stripped) > 1:
            flush_replica()
            current_speaker = stripped[:-1].strip()
            continue

        # порожній рядок
        if stripped == "":
            if current_speaker is not None:
                # просто перенос рядка всередині репліки
                current_buffer.append("")
            continue

        # текст усередині репліки
        if current_speaker is not None:
            current_buffer.append(stripped)
        else:
            # текст поза будь-якою реплікою = ремарка між репліками
            remark_blocks.append(stripped)

    flush_replica()
    return replicas, remark_blocks


def split_replica_by_sentences(text: str, max_chars: int):
    """Ділить одну репліку на кілька частин по розділових знаках."""
    parts = re.split(r"([\.!\?;:])", text)
    chunks: list[str] = []
    buf = ""

    for i in range(0, len(parts), 2):
        part = parts[i].strip()
        end = parts[i + 1] if i + 1 < len(parts) else ""
        if not part and not end:
            continue

        sentence = (part + end).strip()
        if not sentence:
            continue

        candidate = (buf + " " + sentence).strip() if buf else sentence
        if len(candidate) <= max_chars:
            buf = candidate
        else:
            if buf:
                chunks.append(buf)
            buf = sentence

    if buf:
        chunks.append(buf)

    return chunks



def script_file_to_chunks(raw_text: str, settings) -> list[str]:
    max_chars = get_max_chars_for_font_size(settings["font_size"])
    replicas, remark_blocks = parse_replicas_from_colon_headers(raw_text)
    chunks: list[str] = []

    # репліки
    for speaker, rep_text in replicas:
        sentence_chunks = split_replica_by_sentences(rep_text, max_chars=max_chars)
        for part_idx, part in enumerate(sentence_chunks):
            if part_idx == 0:
                # Перша частина репліки: ім'я + двокрапка + текст
                chunks.append(f"{speaker}: \n{part}")
            else:
                # Наступні частини: тільки текст
                chunks.append(part)

    # ремарки між репліками як окремі слайди
    for remark in remark_blocks:
        chunks.append(remark)

    return chunks


# ================== ДОПОМОЖНІ ФУНКЦІЇ ==================

def get_user_settings(context: ContextTypes.DEFAULT_TYPE):
    if "settings" not in context.user_data:
        context.user_data["settings"] = DEFAULT_SETTINGS.copy()
    return context.user_data["settings"]

def build_presentation(chunks, settings) -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    font_name = settings["font_name"]
    font_size = settings["font_size"]
    text_color = settings.get("text_color")

    for chunk_idx, chunk in enumerate(chunks):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # фон
        if settings["bg_mode"] == "color":
            fill = slide.background.fill
            fill.solid()
            r, g, b = settings["bg_color"]
            fill.fore_color.rgb = RGBColor(r, g, b)
        elif settings["bg_mode"] == "image" and settings["bg_image"]:
            img_stream = BytesIO(settings["bg_image"].getvalue())
            slide.shapes.add_picture(
                img_stream,
                left=Inches(0),
                top=Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

        # текстбокс — ВАЖЛИВО: використовуємо TEXT_AREA_* і включаємо перенесення
        left = Inches(TEXT_AREA_LEFT_IN)
        top = Inches(TEXT_AREA_TOP_IN)
        width = Inches(TEXT_AREA_WIDTH_IN)
        height = Inches(TEXT_AREA_HEIGHT_IN)
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.word_wrap = True  # ← тепер текст переноситься по ширині

        # розбити chunk на ім'я та текст
        if ":\n" in chunk:
            speaker_line, body_text = chunk.split(":\n", 1)
        else:
            speaker_line = None
            body_text = chunk

        # Визначити, чи це перша частина репліки персонажа
        is_first_chunk_of_speaker = True
        if chunk_idx > 0:
            prev_chunk = chunks[chunk_idx - 1]
            if ":\n" in prev_chunk:
                prev_speaker = prev_chunk.split(":\n", 1)[0]
                if prev_speaker == speaker_line:
                    is_first_chunk_of_speaker = False

        if speaker_line and is_first_chunk_of_speaker:
            # 1) ІМ'Я ПЕРСОНАЖА ЗЛІВА з двокрапкою
            p_name = tf.paragraphs[0]
            p_name.alignment = PP_ALIGN.LEFT
            run_name = p_name.add_run()
            run_name.text = speaker_line + ":"
            font_name_run = run_name.font
            font_name_run.size = Pt(font_size)
            font_name_run.name = font_name
            font_name_run.bold = True
            if text_color:
                font_name_run.color.rgb = RGBColor(*text_color)

            # 2) ТЕКСТ РЕПЛІКИ ПО ЦЕНТРУ
            if body_text:
                p_body = tf.add_paragraph()
                p_body.alignment = PP_ALIGN.CENTER
                run_body = p_body.add_run()
                run_body.text = body_text
                font_body = run_body.font
                font_body.size = Pt(font_size)
                font_body.name = font_name
                font_body.bold = True
                if text_color:
                    font_body.color.rgb = RGBColor(*text_color)

        else:
            # 3) ПРОДОВЖЕННЯ РЕПЛІКИ АБО РЕМАРКА — ВСЕ ПО ЦЕНТРУ
            p_body = tf.paragraphs[0]
            p_body.alignment = PP_ALIGN.CENTER
            run_body = p_body.add_run()
            run_body.text = body_text
            font_body = run_body.font
            font_body.size = Pt(font_size)
            font_body.name = font_name
            font_body.bold = True
            if text_color:
                font_body.color.rgb = RGBColor(*text_color)

    bio = BytesIO()
    bio.name = "presentation.pptx"
    prs.save(bio)
    bio.seek(0)
    return bio


# ================== КНОПКИ НАЛАШТУВАНЬ ==================

def build_main_keyboard():
    return InlineKeyboardMarkup([
        [
            InlineKeyboardButton("Шрифт", callback_data="cfg_font"),
            InlineKeyboardButton("Розмір", callback_data="cfg_size"),
        ],
        [
            InlineKeyboardButton("Фон", callback_data="cfg_bg"),
        ],
        [
            InlineKeyboardButton("Поточні налаштування", callback_data="cfg_show"),
        ],
    ])

def build_font_keyboard():
    rows = []
    row = []
    for name in AVAILABLE_FONTS:
        row.append(InlineKeyboardButton(name, callback_data=f"font_{name}"))
        if len(row) == 2:
            rows.append(row)
            row = []
    if row:
        rows.append(row)
    rows.append([InlineKeyboardButton("⬅ Назад", callback_data="back_main")])
    return InlineKeyboardMarkup(rows)

def build_size_keyboard():
    rows = []
    row = []
    for sz in AVAILABLE_FONT_SIZES:
        row.append(InlineKeyboardButton(str(sz), callback_data=f"size_{sz}"))
        if len(row) == 3:
            rows.append(row)
            row = []
    if row:
        rows.append(row)
    rows.append([InlineKeyboardButton("⬅ Назад", callback_data="back_main")])
    return InlineKeyboardMarkup(rows)

def build_bg_keyboard():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("Без фону (білий)", callback_data="bg_none")],
        [
            InlineKeyboardButton("Колір: чорний", callback_data="bg_color_0_0_0"),
            InlineKeyboardButton("Колір: сірий", callback_data="bg_color_60_60_60"),
        ],
        [
            InlineKeyboardButton("Колір: синій", callback_data="bg_color_10_20_80"),
            InlineKeyboardButton("Колір: світлий", callback_data="bg_color_240_240_240"),
        ],
        [InlineKeyboardButton("Картинка фону", callback_data="bg_image")],
        [InlineKeyboardButton("⬅ Назад", callback_data="back_main")],
    ])

def build_text_color_keyboard():
    return InlineKeyboardMarkup([
        [
            InlineKeyboardButton("Білий", callback_data="textcol_255_255_255"),
            InlineKeyboardButton("Жовтий", callback_data="textcol_255_255_0"),
        ],
        [
            InlineKeyboardButton("Червоний", callback_data="textcol_255_0_0"),
            InlineKeyboardButton("Зелений", callback_data="textcol_0_255_0"),
        ],
        [InlineKeyboardButton("⬅ Назад", callback_data="back_main")],
    ])


async def settings_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    data = query.data
    settings = get_user_settings(context)
    await query.answer()

    if data == "cfg_font":
        await query.edit_message_text(
            "Оберіть шрифт:",
            reply_markup=build_font_keyboard()
        )

    elif data.startswith("font_"):
        font_name = data.split("font_", 1)[1]
        if font_name in AVAILABLE_FONTS:
            settings["font_name"] = font_name
        await query.edit_message_text(
            f"Шрифт: {settings['font_name']}",
            reply_markup=build_main_keyboard()
        )

    elif data == "cfg_size":
        await query.edit_message_text(
            "Оберіть розмір шрифту:",
            reply_markup=build_size_keyboard()
        )

    elif data.startswith("size_"):
        size = int(data.split("size_", 1)[1])
        if size in AVAILABLE_FONT_SIZES:
            settings["font_size"] = size
        await query.edit_message_text(
            f"Розмір шрифту: {settings['font_size']}",
            reply_markup=build_main_keyboard()
        )

    elif data == "cfg_bg":
        await query.edit_message_text(
            "Налаштування фону:",
            reply_markup=build_bg_keyboard()
        )

    elif data == "bg_none":
        settings["bg_mode"] = "none"
        settings["bg_image"] = None
        await query.edit_message_text(
            "Фон: без фону (білий).",
            reply_markup=build_main_keyboard()
        )

    elif data.startswith("bg_color_"):
        parts = data.split("_")  # ['bg','color','r','g','b']
        r, g, b = int(parts[2]), int(parts[3]), int(parts[4])
        settings["bg_mode"] = "color"
        settings["bg_color"] = (r, g, b)
        settings["bg_image"] = None

        brightness = 0.299 * r + 0.587 * g + 0.114 * b
        if brightness < 128:
            await query.edit_message_text(
                "Темний фон. Оберіть колір тексту:",
                reply_markup=build_text_color_keyboard()
            )
        else:
            await query.edit_message_text(
                f"Фон: колір {settings['bg_color']}",
                reply_markup=build_main_keyboard()
            )

    elif data.startswith("textcol_"):
        parts = data.split("_")  # ['textcol','r','g','b']
        r, g, b = int(parts[1]), int(parts[2]), int(parts[3])
        settings["text_color"] = (r, g, b)
        await query.edit_message_text(
            f"Колір тексту: {settings['text_color']}",
            reply_markup=build_main_keyboard()
        )

    elif data == "bg_image":
        settings["await_bg_image"] = True
        settings["bg_mode"] = "image"
        await query.edit_message_text(
            "Надішли мені картинку – вона стане фоном для всієї презентації.\n"
            "Після цього можна буде надіслати текст або файл сценарію.",
            reply_markup=build_main_keyboard()
        )

    elif data == "cfg_show":
        await query.edit_message_text(
            "Поточні налаштування:\n"
            f"- Шрифт: {settings['font_name']}\n"
            f"- Розмір: {settings['font_size']}\n"
            f"- Режим фону: {settings['bg_mode']}\n"
            f"- Колір фону: {settings['bg_color']}",
            reply_markup=build_main_keyboard()
        )

    elif data == "back_main":
        await query.edit_message_text(
            "Налаштування презентації:",
            reply_markup=build_main_keyboard()
        )

# ================== НАКОПИЧЕННЯ ТЕКСТУ ==================

def get_accumulated_text(context: ContextTypes.DEFAULT_TYPE) -> str:
    """Повертає накопичений текст."""
    return context.user_data.get("accumulated_text", "")

def add_to_accumulated_text(context: ContextTypes.DEFAULT_TYPE, text: str):
    """Додає текст до накопиченого."""
    current = get_accumulated_text(context)
    if current:
        context.user_data["accumulated_text"] = current + "\n" + text
    else:
        context.user_data["accumulated_text"] = text

def clear_accumulated_text(context: ContextTypes.DEFAULT_TYPE):
    """Очищає накопичений текст."""
    context.user_data["accumulated_text"] = ""
async def generate_presentation_from_accumulated(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Генерує презентацію з накопиченого тексту."""
    settings = get_user_settings(context)
    accumulated = get_accumulated_text(context)
    
    if not accumulated:
        return
    
    # Використовуємо логіку сценарію (як для файлів)
    chunks = script_file_to_chunks(accumulated, settings)
    
    if not chunks:
        await update.message.reply_text("Не вдалося обробити текст.")
        clear_accumulated_text(context)
        return
    
    # Назва файлу з першого рядка тексту
    first_line = accumulated.split('\n')[0].strip()
    
    # Якщо перший рядок закінчується на ":" (це ім'я персонажа), беремо його без двокрапки
    if first_line.endswith(':'):
        filename = first_line[:-1].strip()
    else:
        filename = first_line
    
    # Якщо перший рядок порожній або дуже короткий
    if not filename or len(filename) < 3:
        filename = "presentation"
    
    # Очищаємо назву від недозволених символів та обмежуємо довжину
    filename = re.sub(r'[<>:"/\\|?*]', '', filename)
    filename = filename[:50]
    filename = filename + ".pptx"
    
    await update.message.reply_text(f"Генерую презентацію ({len(chunks)} слайдів)...")
    
    pptx_file = build_presentation(chunks, settings)
    await update.message.reply_document(
        document=InputFile(pptx_file, filename=filename),
        caption=f"Презентація готова!\nСлайдів: {len(chunks)}"
    )
    
    # Очищаємо після генерації
    clear_accumulated_text(context)
    # Скасовуємо таймер, якщо він був
    if "timer_task" in context.user_data:
        context.user_data["timer_task"].cancel()
        context.user_data.pop("timer_task", None)


# ================== ХЕНДЛЕРИ БОТА ==================

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    get_user_settings(context)
    clear_accumulated_text(context)
    await update.message.reply_text(
        "Привіт!\n\n"
        "1) Налаштуй шрифт / розмір / фон через кнопки.\n"
        "2) Надішли текст (можна кілька повідомлень підряд).\n"
        "3) Через 3 секунди після останнього повідомлення автоматично створюється презентація.\n"
        "4) Або надішли .txt файл — одразу згенерується презентація.\n\n"
        "Формат тексту:\n"
        "Персонаж 1:\n"
        "текст репліки\n\n"
        "Персонаж 2:\n"
        "текст репліки\n",
        reply_markup=build_main_keyboard()
    )

async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    settings = get_user_settings(context)
    if not settings.get("await_bg_image"):
        return
    photo = update.message.photo[-1]
    file = await photo.get_file()
    bio = BytesIO()
    await file.download_to_memory(out=bio)
    bio.seek(0)
    settings["bg_image"] = bio
    settings["await_bg_image"] = False
    await update.message.reply_text(
        "Фон-картинку збережено. Вона буде застосована до всіх слайдів.\n"
        "Тепер можеш надіслати текст або .txt сценарій.",
        reply_markup=build_main_keyboard()
    )

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Сценарний режим для .txt файлів."""
    settings = get_user_settings(context)
    doc = update.message.document
    if not doc.file_name.lower().endswith(".txt"):
        await update.message.reply_text("Поки що для сценарію підтримую тільки .txt файли.")
        return

    file = await doc.get_file()
    bio = BytesIO()
    await file.download_to_memory(out=bio)
    bio.seek(0)
    raw_bytes = bio.read()
    try:
        text = raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = raw_bytes.decode("cp1251")
        except UnicodeDecodeError:
            await update.message.reply_text("Не вдалося прочитати файл як текст (UTF‑8/CP1251).")
            return

    chunks = script_file_to_chunks(text, settings)
    if not chunks:
        await update.message.reply_text("Не вдалося знайти репліки в цьому сценарії.")
        return

    # Назва файлу — з імені завантаженого файлу (без .txt)
    original_filename = doc.file_name
    if original_filename.lower().endswith('.txt'):
        filename = original_filename[:-4]  # прибираємо .txt
    else:
        filename = original_filename
    
    # Очищаємо назву від недозволених символів
    filename = re.sub(r'[<>:"/\\|?*]', '', filename)
    filename = filename[:50]  # обмежуємо до 50 символів
    filename = filename + ".pptx"

    await update.message.reply_text("Генерую презентацію зі сценарію, зачекай...")
    pptx_file = build_presentation(chunks, settings)
    await update.message.reply_document(
        document=InputFile(pptx_file, filename=filename),
        caption=f"Ось презентація за сценарієм.\nСлайдів: {len(chunks)}"
    )


async def timer_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Викликається через TEXT_ACCUMULATION_TIMEOUT після останнього повідомлення."""
    await generate_presentation_from_accumulated(update, context)

async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Накопичує текст і запускає таймер."""
    text = update.message.text.strip()
    
    # Додаємо текст до накопиченого
    add_to_accumulated_text(context, text)
    
    # Скасовуємо попередній таймер, якщо він був
    if "timer_task" in context.user_data:
        context.user_data["timer_task"].cancel()
    
    # Створюємо новий таймер
    async def delayed_generation():
        await asyncio.sleep(TEXT_ACCUMULATION_TIMEOUT)
        await generate_presentation_from_accumulated(update, context)
    
    task = asyncio.create_task(delayed_generation())
    context.user_data["timer_task"] = task
    
    accumulated = get_accumulated_text(context)
    await update.message.reply_text(
        f"✅ Текст додано ({len(accumulated)} символів)\n"
        f"Якщо це все — чекай {TEXT_ACCUMULATION_TIMEOUT} сек, презентація згенерується автоматично.\n"
        f"Або надішли ще текст."
    )

def main():
    if not BOT_TOKEN or BOT_TOKEN.startswith("ВСТАВ_СВІЙ_ТОКЕН"):
        raise RuntimeError("Спочатку впиши BOT_TOKEN у коді.")
    app = ApplicationBuilder().token(BOT_TOKEN).build()
    app.add_handler(CommandHandler("start", start))
    app.add_handler(CallbackQueryHandler(settings_callback))
    app.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))
    app.run_polling()

if __name__ == "__main__":
    main()
